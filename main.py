import os
from pptx import Presentation
from pptx.util import Inches
from docx import Document
import win32com.client  # For adding sections and Section Zoom

def extract_headings(doc_path):
    doc = Document(doc_path)
    headings = []
    current_section = None
    
    for para in doc.paragraphs:
        if para.style.name.startswith("Heading 1"):
            current_section = {"title": para.text, "slides": []}
            headings.append(current_section)
        elif para.style.name.startswith("Heading 2") and current_section:
            slide_content = {"title": para.text, "text": ""}
            current_section["slides"].append(slide_content)
        elif para.style.name.startswith("Heading 3") or para.style.name.startswith("Heading 4"):
            if current_section and current_section["slides"]:
                current_section["slides"][-1]["text"] += f"\n{para.text}"
    
    return headings

def create_presentation(headings, output_ppt):
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Automatische PowerPoint"
    subtitle.text = "Gegenereerd uit Word-document"

    # Summary Slide
    summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
    summary_slide.shapes.title.text = "Overzicht"

    slide_numbers = []
    for section in headings:
        section_slide = prs.slides.add_slide(prs.slide_layouts[1])
        section_slide.shapes.title.text = section["title"]
        slide_numbers.append(len(prs.slides))  # Store slide index for section zoom

        for slide_content in section["slides"]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_content["title"]
            slide.placeholders[1].text = slide_content["text"]

    # Thank You Slide
    thank_you_slide = prs.slides.add_slide(prs.slide_layouts[1])
    thank_you_slide.shapes.title.text = "Bedankt voor uw aandacht!"
    thank_you_slide.placeholders[1].text = "Vragen? Neem contact op."

    prs.save(output_ppt)
    print(f"PowerPoint opgeslagen als {output_ppt}")

    return slide_numbers  # Return slide indexes for section zoom

def add_sections_and_zoom(ppt_path, section_titles, slide_numbers):
    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    ppt_app.Visible = True  # Show PowerPoint while processing

    presentation = ppt_app.Presentations.Open(os.path.abspath(ppt_path))

    # Add sections
    for i, title in enumerate(section_titles):
        slide_index = slide_numbers[i]
        try:
            presentation.SectionProperties.AddBeforeSlide(slide_index, title)
        except Exception as e:
            print(f"Error adding section: {e}")

    # Save and close
    presentation.Save()
    presentation.Close()
    ppt_app.Quit()
    print("Sections added successfully.")

# Example usage
doc_path = "input.docx"
output_ppt = "output.pptx"

headings = extract_headings(doc_path)
slide_numbers = create_presentation(headings, output_ppt)

# Extract section titles from headings
section_titles = [section["title"] for section in headings]

add_sections_and_zoom(output_ppt, section_titles, slide_numbers)
